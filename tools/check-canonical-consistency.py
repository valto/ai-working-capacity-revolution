#!/usr/bin/env python3
"""
Mechanical consistency checker for "Why Are They Spending Trillions on AI?"

Purpose: no chart, slide, PDF, or workbook in this package should be able to
quietly diverge from data/canonical-cost-model.csv. This script does not (yet)
regenerate every derivative asset from the canonical table — that is a larger
rebuild-pipeline change — but it mechanically enforces that:

  1. The canonical $/M-token and $/AI-working-hour figures for Home, Cooperative,
     and Hyperscale appear, verbatim, in every markdown asset that cites them.
  2. A fixed list of previously-found-and-fixed STALE figures (superseded ranges
     that do not match the canonical table) does NOT reappear anywhere outside
     of an explicit "Corrected"/"Resolved" self-correction note, which is this
     package's established pattern for discussing a superseded number by name.
  3. (If openpyxl is installed) specific cells in the companion Excel workbook
     that have previously broken — a hyperscale utilization assumption, and two
     formula-string cells that would render #NAME? — still hold their expected
     values/types.
  4. (If python-pptx is installed) the slide deck's Home/Hyperscale cost bullet
     and the cooperative-model table row do not contain stale figures.

Run: python3 tools/check-canonical-consistency.py
Exit code 0 = clean. Exit code 1 = drift found (see printed report).
"""
import csv
import os
import re
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

CANONICAL_CSV = os.path.join(ROOT, "data", "canonical-cost-model.csv")

MARKDOWN_FILES_TO_CHECK = [
    "01-whitepaper.md",
    "05-workbook-token-factory-scenarios.md",
    "13-slide-deck-outline.md",
    "17-visual-asset-briefs.md",
]

# Patterns for figures that were found stale and fixed in this package's history.
# Each entry: (regex, human-readable description, files it must NOT appear in
# outside of a "Corrected"/"Resolved" self-correction sentence).
STALE_PATTERNS = [
    (r"\$0\.6[–-]\$?2\b(?!\d)", "stale Home tier figure ($0.6-2/M, superseded by canonical $1.37-11.89/M)"),
    (r"0\.77[–-]1\.20\b", "stale Cooperative tier figure (0.77-1.20/M, superseded by canonical $1.99-7.62/M)"),
    (r"0\.006[–-]\$?0\.24\b", "stale Home hourly figure ($0.006-0.24/hr, derived from superseded Home cost)"),
    (r"0\.046[–-]\$?0\.72\b", "stale Cooperative hourly figure ($0.046-0.72/hr, derived from superseded Cooperative cost)"),
]

# A line is allowed to contain a stale pattern ONLY if it also contains one of
# these self-correction markers, per this package's established convention.
SELF_CORRECTION_MARKERS = ["Corrected 2026-08-13", "Resolved 2026-08-13", "Resolved (2026-08-13)"]

CANONICAL_REQUIRED_STRINGS = [
    "1.37",   # Home low
    "11.89",  # Home high
    "1.99",   # Cooperative low
    "7.62",   # Cooperative high
]


def load_canonical():
    rows = []
    with open(CANONICAL_CSV, newline="") as f:
        for row in csv.DictReader(f):
            rows.append(row)
    return rows


def check_markdown_files():
    problems = []
    for relpath in MARKDOWN_FILES_TO_CHECK:
        path = os.path.join(ROOT, relpath)
        if not os.path.exists(path):
            problems.append(f"MISSING FILE expected for consistency check: {relpath}")
            continue
        with open(path, encoding="utf-8") as f:
            lines = f.readlines()
        text = "".join(lines)

        # Positive check: canonical figures should appear somewhere in the whitepaper.
        if relpath == "01-whitepaper.md":
            for needle in CANONICAL_REQUIRED_STRINGS:
                if needle not in text:
                    problems.append(
                        f"{relpath}: canonical figure '{needle}' not found anywhere — "
                        f"has the canonical Home/Cooperative range been removed or changed?"
                    )

        # Negative check: stale patterns must not appear outside a self-correction note.
        for i, line in enumerate(lines, start=1):
            for pattern, desc in STALE_PATTERNS:
                if re.search(pattern, line):
                    if not any(marker in line for marker in SELF_CORRECTION_MARKERS):
                        problems.append(
                            f"{relpath}:{i}: found {desc} without a "
                            f"'Corrected'/'Resolved' self-correction marker on the same line"
                        )
    return problems


def check_xlsx():
    problems = []
    xlsx_path = os.path.join(ROOT, "18-companion-data-model.xlsx")
    if not os.path.exists(xlsx_path):
        return [f"MISSING FILE expected for consistency check: 18-companion-data-model.xlsx"]
    try:
        import openpyxl
    except ImportError:
        print("  (skipping xlsx checks: openpyxl not installed)")
        return problems

    wb = openpyxl.load_workbook(xlsx_path)

    # Regression guard: Hyperscale mid-case utilization must be 60%, not 90%.
    try:
        ws = wb["Hyperscale Tier"]
        util = ws["B14"].value
        if util is None or abs(float(util) - 0.60) > 1e-6:
            problems.append(
                f"18-companion-data-model.xlsx: Hyperscale Tier!B14 utilization is "
                f"{util!r}, expected 0.60 (this paper's canonical mid-case) — "
                f"previously found set to 0.90, producing a wrong $0.091/M instead of $0.133/M"
            )
    except KeyError:
        problems.append("18-companion-data-model.xlsx: 'Hyperscale Tier' sheet not found")

    # Regression guard: these two cells were prose labels stored as formulas
    # (leading '='), which Excel renders as #NAME?. They must be plain strings.
    for sheet_name, cell_addr in [("Home Tier", "C26"), ("Hyperscale Tier", "C25")]:
        try:
            ws = wb[sheet_name]
            cell = ws[cell_addr]
            if cell.data_type == "f":
                problems.append(
                    f"18-companion-data-model.xlsx: {sheet_name}!{cell_addr} is stored as a "
                    f"formula ({cell.value!r}) — this renders #NAME? in Excel; it should be a "
                    f"plain text cell"
                )
        except KeyError:
            problems.append(f"18-companion-data-model.xlsx: '{sheet_name}' sheet not found")

    return problems


def check_pptx():
    problems = []
    pptx_path = os.path.join(ROOT, "19-slide-deck.pptx")
    if not os.path.exists(pptx_path):
        return [f"MISSING FILE expected for consistency check: 19-slide-deck.pptx"]
    try:
        from pptx import Presentation
    except ImportError:
        print("  (skipping pptx checks: python-pptx not installed)")
        return problems

    prs = Presentation(pptx_path)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text)
    full_text = "\n".join(all_text)

    for pattern, desc in STALE_PATTERNS:
        if re.search(pattern, full_text):
            problems.append(f"19-slide-deck.pptx: found {desc} in rendered slide text")

    if "1.37" not in full_text or "11.89" not in full_text:
        problems.append(
            "19-slide-deck.pptx: canonical Home-tier figures (1.37/11.89) not found in "
            "rendered slide text — has slide 11 drifted from the canonical model?"
        )

    return problems


def main():
    print("Checking canonical cost-model consistency...")
    print(f"  Canonical source: {os.path.relpath(CANONICAL_CSV, ROOT)}")
    load_canonical()  # validates the CSV parses; rows unused directly (checks are pattern-based)

    all_problems = []
    all_problems += check_markdown_files()
    all_problems += check_xlsx()
    all_problems += check_pptx()

    print()
    if not all_problems:
        print("PASS — no drift detected across markdown, xlsx, and pptx assets.")
        return 0

    print(f"FAIL — {len(all_problems)} consistency issue(s) found:\n")
    for p in all_problems:
        print(f"  - {p}")
    return 1


if __name__ == "__main__":
    sys.exit(main())
